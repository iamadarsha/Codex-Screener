#!/usr/bin/env python3
"""
Codex Screener — Business Pitch Deck Generator
Creates a professional 16-slide dark-theme PPTX presentation.
Owner: Adarsha Chatterjee
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0a, 0x0e, 0x1a)
BG_CARD    = RGBColor(0x11, 0x18, 0x27)
BG_CARD2   = RGBColor(0x1e, 0x29, 0x3b)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x9c, 0xa3, 0xaf)
EMERALD    = RGBColor(0x10, 0xb9, 0x81)
PURPLE     = RGBColor(0x8b, 0x5c, 0xf6)
CYAN       = RGBColor(0x06, 0xb6, 0xd4)
ORANGE     = RGBColor(0xf9, 0x73, 0x16)
RED        = RGBColor(0xef, 0x44, 0x44)
BLUE       = RGBColor(0x3b, 0x82, 0xf6)
YELLOW     = RGBColor(0xfb, 0xbf, 0x24)
AMBER      = RGBColor(0xf5, 0x9e, 0x0b)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_bar(slide, left, top, width, height, color=EMERALD):
    return add_rect(slide, left, top, width, height, color)

def add_stat_card(slide, left, top, width, height, number, label, accent=EMERALD):
    card = add_shape(slide, left, top, width, height, BG_CARD)
    add_accent_bar(slide, left, top, width, Inches(0.05), accent)
    add_text(slide, left, top + Inches(0.25), width, Inches(0.6), number, 36, accent, True, PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.85), width, Inches(0.4), label, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    return card

def slide_title(slide, text, subtitle=None):
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), text, 36, WHITE, True)
    add_accent_bar(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.06), EMERALD)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5), subtitle, 16, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), EMERALD)
add_rect(s, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), PURPLE)

add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "CODEX SCREENER", 60, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
         "AI-Powered Indian Stock Market Intelligence Platform", 28, EMERALD, False, PP_ALIGN.CENTER)
add_accent_bar(s, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06), PURPLE)
add_text(s, Inches(1), Inches(3.9), Inches(11), Inches(0.5),
         "Real-time screener  |  AI stock picks  |  500+ NSE stocks  |  13 technical scans",
         16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "Founded by Adarsha Chatterjee", 20, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
         "adarsha.chatterjee@gmail.com  |  March 2026", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "The Problem", "Why Indian retail investors struggle")

problems = [
    ("Fragmented Tools", "Investors juggle 5-7 apps for screeners,\ncharts, news, and alerts — losing time\nand missing trades.", RED),
    ("No AI for NSE", "Global AI tools focus on US markets.\nIndian investors have zero AI-powered\nstock analysis built for NSE.", ORANGE),
    ("Expensive Terminals", "Bloomberg ($24K/yr) and Refinitiv\n($22K/yr) are out of reach for 100M+\nIndian retail investors.", YELLOW),
    ("Information Overload", "10+ financial news sources, 500+ stocks,\ndozens of indicators — no unified\nintelligent layer.", PURPLE),
]

x = Inches(0.6)
for title, desc, color in problems:
    card = add_shape(s, x, Inches(2.0), Inches(2.9), Inches(4.5), BG_CARD, color)
    add_accent_bar(s, x, Inches(2.0), Inches(2.9), Inches(0.06), color)
    add_text(s, x + Inches(0.2), Inches(2.3), Inches(2.5), Inches(0.5), title, 20, color, True)
    add_text(s, x + Inches(0.2), Inches(3.0), Inches(2.5), Inches(3.0), desc, 14, LIGHT_GRAY)
    x += Inches(3.15)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Our Solution", "One platform. All intelligence. Zero cost to start.")

solutions = [
    ("AI Stock Picker", "3-layer AI engine analyzes 500+ stocks\nacross 10 news feeds + technical\nindicators. Picks 5 stocks per\ntimeframe daily.", EMERALD),
    ("Smart Screener", "13 prebuilt institutional-grade scans:\nGolden Cross, RSI Oversold, Bollinger\nSqueeze, Volume Breakout, and more.", PURPLE),
    ("Real-Time Data", "Live prices for Nifty 500 via NSE\npolling every 30 seconds. WebSocket\nstreaming. Redis-cached for\nsub-100ms response.", CYAN),
    ("Full Terminal", "Charts (TradingView), Watchlists,\nAlerts, Fundamentals, Company Info —\neverything in one dark-themed\nprofessional UI.", BLUE),
]

x = Inches(0.6)
for title, desc, color in solutions:
    card = add_shape(s, x, Inches(2.0), Inches(2.9), Inches(4.5), BG_CARD, color)
    add_accent_bar(s, x, Inches(2.0), Inches(2.9), Inches(0.06), color)
    add_text(s, x + Inches(0.2), Inches(2.3), Inches(2.5), Inches(0.5), title, 20, color, True)
    add_text(s, x + Inches(0.2), Inches(3.0), Inches(2.5), Inches(3.0), desc, 14, LIGHT_GRAY)
    x += Inches(3.15)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — How It Works
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "How It Works", "AI-first architecture with intelligent fallback")

steps = [
    ("1", "DATA INGESTION", "NSE Poller fetches live prices\nfor 500 stocks every 30s.\nYahoo Finance provides\nhistorical OHLCV data.", EMERALD),
    ("2", "INDICATOR ENGINE", "Computes RSI, EMA, MACD,\nBollinger Bands, ATR, OBV\nfor all 500 stocks.\nCached in Redis.", PURPLE),
    ("3", "AI ANALYSIS", "3-Layer engine:\nL1: RSS + Technical (free)\nL2: Gemini 3.1 Flash Lite\nL3: Groq Llama 3.3", CYAN),
    ("4", "SMART DELIVERY", "Results cached until next\nmarket open. Instant API.\nWebSocket for real-time\nupdates.", ORANGE),
]

x = Inches(0.4)
for num, title, desc, color in steps:
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), Inches(2.0), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if num != "4":
        add_rect(s, x + Inches(2.95), Inches(2.15), Inches(0.3), Inches(0.06), color)

    add_text(s, x + Inches(0.2), Inches(2.9), Inches(2.8), Inches(0.4), title, 16, color, True, PP_ALIGN.CENTER)
    card = add_shape(s, x + Inches(0.15), Inches(3.4), Inches(2.8), Inches(3.0), BG_CARD)
    add_text(s, x + Inches(0.35), Inches(3.6), Inches(2.4), Inches(2.6), desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    x += Inches(3.2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — AI Stock Picker Deep Dive
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "AI Stock Picker", "3-layer intelligent fallback — always delivers picks")

layers = [
    ("LAYER 1: RSS + TECHNICAL", "Zero API Cost", [
        "Parses 10 RSS news feeds (Google News, Moneycontrol, ET, Mint, BS)",
        "Scores 500 stocks using technical indicators (RSI, EMA, MACD)",
        "Combines sentiment + momentum for composite score",
        "Picks top 5 per timeframe: Intraday, Weekly, Monthly",
    ], EMERALD, "$0/day"),
    ("LAYER 2: GEMINI AI", "Google Gemini 3.1 Flash Lite", [
        "Sends market context + news summary to Gemini",
        "AI generates structured JSON stock recommendations",
        "Includes confidence scores and reasoning",
        "500 requests/day free tier — fully funded",
    ], PURPLE, "Free tier"),
    ("LAYER 3: GROQ + xAI", "Fallback LLM Layer", [
        "Groq Llama 3.3 70B — ultra-fast inference",
        "xAI Grok as final fallback",
        "Same structured output format",
        "Triple redundancy = 99.9% uptime",
    ], CYAN, "Pay-per-use"),
]

y = Inches(1.8)
for title, subtitle, bullets, color, cost in layers:
    card = add_shape(s, Inches(0.6), y, Inches(8.5), Inches(1.6), BG_CARD, color)
    add_accent_bar(s, Inches(0.6), y, Inches(0.08), Inches(1.6), color)
    add_text(s, Inches(1.0), y + Inches(0.15), Inches(4), Inches(0.4), title, 18, color, True)
    add_text(s, Inches(1.0), y + Inches(0.5), Inches(3.5), Inches(0.3), subtitle, 12, LIGHT_GRAY)
    bx = Inches(5.0)
    by = y + Inches(0.15)
    for b in bullets:
        add_text(s, bx, by, Inches(4), Inches(0.3), f"  {b}", 11, LIGHT_GRAY)
        by += Inches(0.3)
    badge = add_shape(s, Inches(9.5), y + Inches(0.5), Inches(1.8), Inches(0.5), color)
    add_text(s, Inches(9.5), y + Inches(0.55), Inches(1.8), Inches(0.4), cost, 14, WHITE, True, PP_ALIGN.CENTER)
    y += Inches(1.8)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Technical Screener
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "13 Prebuilt Scans", "Institutional-grade screening for retail investors")

scans_left = [
    "Golden Cross (50 > 200 EMA)", "Death Cross (50 < 200 EMA)",
    "RSI Oversold (< 30)", "RSI Overbought (> 70)",
    "MACD Bullish Crossover", "Bollinger Squeeze",
    "Volume Breakout (2x avg)",
]
scans_right = [
    "52-Week High Breakout", "52-Week Low Reversal",
    "Bullish Engulfing Pattern", "Bearish Engulfing Pattern",
    "Price Above All EMAs", "Price Below All EMAs",
]

card_l = add_shape(s, Inches(0.6), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD, EMERALD)
add_accent_bar(s, Inches(0.6), Inches(1.8), Inches(5.5), Inches(0.06), EMERALD)
y = Inches(2.1)
for i, scan in enumerate(scans_left):
    c = EMERALD if i % 2 == 0 else CYAN
    add_text(s, Inches(1.0), y, Inches(5), Inches(0.4), f"  {scan}", 15, c)
    y += Inches(0.45)

card_r = add_shape(s, Inches(6.5), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD, PURPLE)
add_accent_bar(s, Inches(6.5), Inches(1.8), Inches(5.5), Inches(0.06), PURPLE)
y = Inches(2.1)
for i, scan in enumerate(scans_right):
    c = PURPLE if i % 2 == 0 else BLUE
    add_text(s, Inches(6.9), y, Inches(5), Inches(0.4), f"  {scan}", 15, c)
    y += Inches(0.45)

add_text(s, Inches(6.9), Inches(5.2), Inches(5), Inches(0.8),
         "All scans run across entire\nNifty 500 universe in < 2 seconds", 14, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Product Features
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Product Features", "Everything a trader needs in one platform")

features = [
    ("Dashboard", "Market overview with stat cards,\nindex ticker bar, breadth donut,\nand breakout feed", EMERALD),
    ("AI Picks", "Daily AI-generated stock picks\nfor 3 timeframes with live\nprice streaming", PURPLE),
    ("Screener", "13 prebuilt scans + custom\nscan builder across Nifty 500\nuniverse", CYAN),
    ("Charts", "TradingView-powered charts\nwith candlestick, volume,\nRSI, MACD, BB, EMA overlays", BLUE),
    ("Watchlist", "Personal stock watchlist with\nadd/remove and real-time\nlive price updates", ORANGE),
    ("Alerts", "Price-based alert triggers\nwith real-time WebSocket\nnotifications", YELLOW),
    ("Fundamentals", "Company financial data,\nratios, and key metrics\nfor informed decisions", EMERALD),
    ("Dark Theme", "Professional terminal-style\nUI with dark/light toggle\nand smooth animations", PURPLE),
]

x, y = Inches(0.4), Inches(1.8)
for i, (title, desc, color) in enumerate(features):
    cx = x + (i % 4) * Inches(3.15)
    cy = y + (i // 4) * Inches(2.6)
    card = add_shape(s, cx, cy, Inches(2.9), Inches(2.3), BG_CARD, color)
    add_accent_bar(s, cx, cy, Inches(2.9), Inches(0.06), color)
    add_text(s, cx + Inches(0.2), cy + Inches(0.2), Inches(2.5), Inches(0.4), title, 18, color, True)
    add_text(s, cx + Inches(0.2), cy + Inches(0.65), Inches(2.5), Inches(1.5), desc, 12, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Tech Stack", "Modern, scalable, production-ready infrastructure")

stacks = [
    ("Frontend", ["Next.js 15 + React 19", "TypeScript 5.9.3", "Tailwind CSS + Framer Motion", "React Query + Zustand", "TradingView Charts"], EMERALD),
    ("Backend", ["FastAPI + Python 3.12", "WebSocket real-time feeds", "SQLAlchemy + Alembic", "APScheduler (30s polling)", "httpx async HTTP client"], PURPLE),
    ("Data Layer", ["Redis (841 cached keys)", "Supabase PostgreSQL", "TimescaleDB hypertables", "OHLCV time-series data", "Nifty 500 seed dataset"], CYAN),
    ("Infrastructure", ["Railway (API + Frontend)", "Docker Compose (local dev)", "Auto-deploy from GitHub", "iOS SwiftUI WebView app", "99.9% uptime target"], BLUE),
]

x = Inches(0.4)
for title, items, color in stacks:
    card = add_shape(s, x, Inches(1.8), Inches(3.0), Inches(4.8), BG_CARD, color)
    add_accent_bar(s, x, Inches(1.8), Inches(3.0), Inches(0.06), color)
    add_text(s, x + Inches(0.3), Inches(2.05), Inches(2.5), Inches(0.5), title, 22, color, True)
    y = Inches(2.6)
    for item in items:
        add_text(s, x + Inches(0.3), y, Inches(2.5), Inches(0.4), f"  {item}", 14, LIGHT_GRAY)
        y += Inches(0.42)
    x += Inches(3.2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Market Opportunity
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Market Opportunity", "India's retail trading revolution")

stats = [
    ("150M+", "Demat Accounts\nin India (2026)", EMERALD),
    ("$4.1T", "NSE Market Cap\n(world's 4th largest)", PURPLE),
    ("42%", "YoY Growth in\nRetail Participation", CYAN),
    ("$2.3B", "Indian WealthTech\nMarket (2026)", ORANGE),
]

x = Inches(0.4)
for num, label, color in stats:
    card = add_shape(s, x, Inches(1.8), Inches(3.0), Inches(2.5), BG_CARD, color)
    add_accent_bar(s, x, Inches(1.8), Inches(3.0), Inches(0.06), color)
    add_text(s, x, Inches(2.1), Inches(3.0), Inches(0.8), num, 44, color, True, PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.95), Inches(3.0), Inches(0.8), label, 15, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    x += Inches(3.2)

insights = [
    ("TAM", "$2.3B Indian WealthTech market growing at 35% CAGR", EMERALD),
    ("SAM", "15M active traders using 3+ trading tools — our direct target", PURPLE),
    ("SOM", "500K users in Year 1 through freemium + organic growth", CYAN),
]

x = Inches(0.6)
for title, desc, color in insights:
    card = add_shape(s, x, Inches(4.8), Inches(3.8), Inches(1.8), BG_CARD, color)
    add_accent_bar(s, x, Inches(4.8), Inches(3.8), Inches(0.06), color)
    add_text(s, x + Inches(0.2), Inches(5.0), Inches(1.5), Inches(0.4), title, 18, color, True)
    add_text(s, x + Inches(0.2), Inches(5.4), Inches(3.4), Inches(1.0), desc, 13, LIGHT_GRAY)
    x += Inches(4.1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Business Model
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Business Model", "Freemium SaaS with multiple revenue streams")

tiers = [
    ("FREE", "\u20b90/mo", ["Dashboard + Market Overview", "5 Prebuilt Scans", "Basic Charts", "1 Watchlist (10 stocks)", "Community Support"], LIGHT_GRAY, EMERALD),
    ("PRO", "\u20b9499/mo", ["All 13 Prebuilt Scans", "AI Stock Picks (3 timeframes)", "Advanced Charts + Indicators", "Unlimited Watchlists", "Price Alerts (50)", "Priority Support"], EMERALD, EMERALD),
    ("PREMIUM", "\u20b9999/mo", ["Everything in Pro", "Custom Scan Builder", "Real-time WebSocket Feeds", "API Access", "Unlimited Alerts", "Dedicated Support"], PURPLE, PURPLE),
    ("INSTITUTIONAL", "Custom", ["White-label Solution", "Bulk API Access", "Custom Integrations", "SLA Guarantee", "Dedicated Infrastructure", "24/7 Support"], CYAN, CYAN),
]

x = Inches(0.3)
for name, price, features, text_color, accent in tiers:
    card = add_shape(s, x, Inches(1.8), Inches(3.05), Inches(5.2), BG_CARD, accent)
    add_accent_bar(s, x, Inches(1.8), Inches(3.05), Inches(0.06), accent)
    add_text(s, x, Inches(2.0), Inches(3.05), Inches(0.5), name, 22, accent, True, PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.5), Inches(3.05), Inches(0.5), price, 28, WHITE, True, PP_ALIGN.CENTER)
    y = Inches(3.2)
    for f in features:
        add_text(s, x + Inches(0.25), y, Inches(2.6), Inches(0.35), f"  {f}", 12, LIGHT_GRAY)
        y += Inches(0.35)
    x += Inches(3.25)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Competitive Advantage
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Competitive Advantage", "What makes Codex Screener different")

advantages = [
    ("AI-First Architecture", "Only platform with a 3-layer AI engine purpose-built for Indian markets.\nCombines news sentiment analysis with technical indicators for superior stock picks.", EMERALD),
    ("Zero-Cost AI Layer", "Layer 1 (RSS + Technical scoring) runs entirely without API calls —\nsustainable unit economics from day one. AI costs scale only with premium users.", PURPLE),
    ("Real-Time Everything", "30-second NSE polling, WebSocket streaming, Redis caching — sub-100ms\nresponse times. Institutional performance at retail prices.", CYAN),
    ("India-Native Design", "Built from ground up for NSE/BSE. All 500 Nifty stocks, INR formatting,\nIST timezone, Indian news sources. Not a US product retrofitted for India.", ORANGE),
]

y = Inches(1.8)
for title, desc, color in advantages:
    card = add_shape(s, Inches(0.6), y, Inches(12), Inches(1.2), BG_CARD, color)
    add_accent_bar(s, Inches(0.6), y, Inches(0.08), Inches(1.2), color)
    add_text(s, Inches(1.0), y + Inches(0.1), Inches(3.5), Inches(0.4), title, 20, color, True)
    add_text(s, Inches(4.5), y + Inches(0.15), Inches(7.5), Inches(1.0), desc, 14, LIGHT_GRAY)
    y += Inches(1.35)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Competitive Landscape
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Competitive Landscape", "How we compare to existing solutions")

headers = ["Feature", "Codex Screener", "Screener.in", "Chartink", "TradingView", "Trendlyne"]
colors_h = [LIGHT_GRAY, EMERALD, LIGHT_GRAY, LIGHT_GRAY, LIGHT_GRAY, LIGHT_GRAY]
widths = [Inches(2.2), Inches(2.2), Inches(2.0), Inches(2.0), Inches(2.2), Inches(2.0)]

x = Inches(0.4)
for i, (h, c, w) in enumerate(zip(headers, colors_h, widths)):
    add_text(s, x, Inches(1.8), w, Inches(0.5), h, 14, c, True, PP_ALIGN.CENTER)
    x += w

rows = [
    ("AI Stock Picks",     ["Yes (3-layer)", "No", "No", "No", "Basic"]),
    ("Technical Screener",  ["13 scans", "Yes", "Yes", "Limited", "Yes"]),
    ("Real-time Prices",    ["30s polling", "Delayed", "Delayed", "Real-time", "15min"]),
    ("NSE/BSE Native",      ["Yes", "Yes", "Yes", "Partial", "Yes"]),
    ("Free Tier",           ["Yes", "Limited", "Limited", "Yes", "Limited"]),
    ("AI News Analysis",    ["10 RSS feeds", "No", "No", "No", "No"]),
    ("WebSocket Streaming", ["Yes", "No", "No", "Yes", "No"]),
    ("Mobile App",          ["iOS (SwiftUI)", "Yes", "No", "Yes", "Yes"]),
    ("Price (Pro)",         ["\u20b9499/mo", "\u20b99K/yr", "\u20b94K/yr", "$15/mo", "\u20b916K/yr"]),
]

y = Inches(2.4)
for idx, (label, vals) in enumerate(rows):
    x = Inches(0.4)
    add_rect(s, Inches(0.3), y - Inches(0.05), Inches(12.7), Inches(0.5), BG_CARD2 if idx % 2 == 0 else BG_CARD)
    add_text(s, x, y, widths[0], Inches(0.4), label, 12, WHITE, True, PP_ALIGN.LEFT)
    x += widths[0]
    for i, v in enumerate(vals):
        c = EMERALD if i == 0 else LIGHT_GRAY
        if v in ("No", "Limited", "Delayed", "Partial"):
            c = RGBColor(0x6b, 0x72, 0x80)
        add_text(s, x, y, widths[i+1], Inches(0.4), v, 11, c, i == 0, PP_ALIGN.CENTER)
        x += widths[i+1]
    y += Inches(0.5)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Traction & Metrics
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Traction & Milestones", "Built and shipped in record time")

milestones = [
    ("Mar 2026", "MVP Launch", "Full platform live on Railway with\nall core features operational", EMERALD),
    ("Mar 2026", "500 Stocks", "Complete Nifty 500 coverage with\nreal-time data pipeline", PURPLE),
    ("Mar 2026", "AI Engine", "3-layer AI stock picker with\n99.9% uptime and fallback chain", CYAN),
    ("Mar 2026", "iOS App", "SwiftUI WebView app ready\nfor App Store submission", BLUE),
]

x = Inches(0.4)
for date, title, desc, color in milestones:
    card = add_shape(s, x, Inches(1.8), Inches(3.0), Inches(2.5), BG_CARD, color)
    add_accent_bar(s, x, Inches(1.8), Inches(3.0), Inches(0.06), color)
    badge = add_shape(s, x + Inches(0.3), Inches(2.0), Inches(1.2), Inches(0.4), color)
    add_text(s, x + Inches(0.3), Inches(2.0), Inches(1.2), Inches(0.4), date, 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), Inches(2.5), Inches(2.5), Inches(0.4), title, 18, WHITE, True)
    add_text(s, x + Inches(0.3), Inches(3.0), Inches(2.5), Inches(1.0), desc, 12, LIGHT_GRAY)
    x += Inches(3.2)

add_text(s, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5), "Technical Stats", 22, EMERALD, True)
stats_row = [
    ("500+", "NSE Stocks", EMERALD), ("841", "Redis Cache Keys", PURPLE),
    ("13", "Prebuilt Scans", CYAN), ("30s", "Polling Cycle", BLUE),
    ("10", "News RSS Feeds", ORANGE), ("<100ms", "API Response", YELLOW),
]
x = Inches(0.4)
for num, label, color in stats_row:
    add_stat_card(s, x, Inches(5.3), Inches(1.95), Inches(1.5), num, label, color)
    x += Inches(2.1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — Go-to-Market Strategy
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "Go-to-Market Strategy", "Land with traders, expand to institutions")

phases = [
    ("Phase 1: Launch", "Q1 2026", [
        "Free tier drives organic signups",
        "Target: Trading communities on Reddit, Twitter/X",
        "YouTube & blog content marketing",
        "Reach: 10K users in 90 days",
    ], EMERALD),
    ("Phase 2: Growth", "Q2-Q3 2026", [
        "Pro tier conversion (target: 5% of free users)",
        "Influencer partnerships (FinTwit India)",
        "SEO-optimized stock pages",
        "Reach: 100K users, 5K paying",
    ], PURPLE),
    ("Phase 3: Scale", "Q4 2026+", [
        "Institutional tier & API marketplace",
        "White-label for brokers & wealth managers",
        "BSE expansion + F&O coverage",
        "Target: 500K users, $1M ARR",
    ], CYAN),
]

x = Inches(0.4)
for title, timeline, bullets, color in phases:
    card = add_shape(s, x, Inches(1.8), Inches(4.0), Inches(5.0), BG_CARD, color)
    add_accent_bar(s, x, Inches(1.8), Inches(4.0), Inches(0.06), color)
    badge = add_shape(s, x + Inches(0.3), Inches(2.0), Inches(1.5), Inches(0.4), color)
    add_text(s, x + Inches(0.3), Inches(2.0), Inches(1.5), Inches(0.4), timeline, 12, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), Inches(2.6), Inches(3.5), Inches(0.5), title, 20, WHITE, True)
    y = Inches(3.2)
    for b in bullets:
        add_text(s, x + Inches(0.4), y, Inches(3.4), Inches(0.45), f"  {b}", 13, LIGHT_GRAY)
        y += Inches(0.45)
    x += Inches(4.2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — The Ask
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "The Ask", "Pre-Seed Round")

add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(0.8),
         "$500,000", 56, EMERALD, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.9), Inches(11), Inches(0.5),
         "Pre-Seed  |  SAFE Note  |  12 Month Runway", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)

uses = [
    ("40%", "Engineering", "2 senior engineers\n+ 1 data scientist", EMERALD),
    ("25%", "Infrastructure", "Railway, Redis, Supabase\nscaling + monitoring", PURPLE),
    ("20%", "Marketing", "Content, influencer\npartnerships, community", CYAN),
    ("15%", "Operations", "Legal, compliance,\nSEBI advisory", ORANGE),
]

x = Inches(0.4)
for pct, title, desc, color in uses:
    card = add_shape(s, x, Inches(3.8), Inches(3.0), Inches(3.0), BG_CARD, color)
    add_accent_bar(s, x, Inches(3.8), Inches(3.0), Inches(0.06), color)
    add_text(s, x, Inches(4.0), Inches(3.0), Inches(0.6), pct, 36, color, True, PP_ALIGN.CENTER)
    add_text(s, x, Inches(4.6), Inches(3.0), Inches(0.4), title, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.1), Inches(3.0), Inches(1.2), desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    x += Inches(3.2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — Thank You / Contact
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), EMERALD)
add_rect(s, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), PURPLE)

add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
         "Thank You", 56, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.6), Inches(11), Inches(0.6),
         "Let's build the Bloomberg Terminal for India's 150M+ retail investors.",
         22, EMERALD, False, PP_ALIGN.CENTER)
add_accent_bar(s, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06), PURPLE)

contact_info = [
    ("Adarsha Chatterjee", 24, WHITE, True),
    ("Founder & CEO", 16, EMERALD, False),
    ("", 10, WHITE, False),
    ("adarsha.chatterjee@gmail.com", 16, LIGHT_GRAY, False),
    ("@iamadarsha", 16, LIGHT_GRAY, False),
    ("", 10, WHITE, False),
    ("Live Demo: breakoutscan-web-production.up.railway.app", 14, CYAN, False),
]

y = Inches(4.0)
for text, size, color, bold in contact_info:
    if text:
        add_text(s, Inches(1), y, Inches(11), Inches(0.5), text, size, color, bold, PP_ALIGN.CENTER)
    y += Inches(0.4)


# ── Save ────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Codex_Screener_Pitch_Deck.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
